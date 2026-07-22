"""Generate 2026-07-28-project-showcase.pptx for the CAMEL Phase I Project Showcase."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# Penn State palette
PSU_BLUE  = RGBColor(0x1E, 0x40, 0x7C)
PSU_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xFA)
ACCENT    = RGBColor(0x96, 0xBE, 0xE6)   # light PSU blue
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY  = RGBColor(0x55, 0x65, 0x70)
FLAG_GOLD = RGBColor(0xE2, 0x97, 0x00)   # for placeholders

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # truly blank


# ── helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, left, top, width, height,
          size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
          wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_lines(tf, lines, size=16, bold=False, color=TEXT_DARK,
              space_before=0, bullet=False):
    """Append multiple paragraph lines to an existing text frame."""
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        if space_before:
            para.space_before = Pt(space_before)
        if bullet:
            para.text = line
        else:
            run = para.add_run()
            run.text = line
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Calibri"


def header_bar(slide, title, subtitle=None):
    """Dark header bar across the top."""
    rect(slide, 0, 0, 13.33, 1.3, fill_rgb=PSU_BLUE)
    txbox(slide, title, 0.4, 0.1, 12.5, 0.65,
          size=28, bold=True, color=PSU_WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.4, 0.78, 12.5, 0.4,
              size=14, color=ACCENT)


def slide_number(slide, n):
    txbox(slide, str(n), 12.8, 7.1, 0.4, 0.3,
          size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ── Slide 1 — Title ──────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=PSU_BLUE)
rect(s, 0, 4.5, 13.33, 3.0, fill_rgb=LIGHT_BG)

txbox(s, "Modeling the Messy", 0.6, 1.0, 12.1, 1.2,
      size=44, bold=True, color=PSU_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Capturing Student Reasoning within\nHigh-Dimensional Industrial AI Datasets",
      0.6, 2.3, 12.1, 1.1,
      size=22, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(s, "NSF CAMEL Phase I  |  Award #2621173  |  Penn State University",
      0.6, 3.6, 12.1, 0.5,
      size=14, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(s, "Napolitano · Yao · Reinhart · Zuo · Hill",
      0.6, 4.85, 12.1, 0.5,
      size=20, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
txbox(s, "July 28, 2026",
      0.6, 5.5, 12.1, 0.4,
      size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 2 — Team Members ───────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
header_bar(s, "Team Members")
slide_number(s, 2)

rows = [
    ("Rebecca (Becca) Napolitano (PI)", "Penn State — AE", "Data translation; Python/notebook education; learning analytics"),
    ("Xiangquan (James) Yao (Co-PI)",   "Penn State",       "Science of learning; cognitive science; educational measurement"),
    ("Wesley F. Reinhart (Co-PI)",      "Penn State — MSE", "FAIR data; scientific AI; Composable Courseware architecture"),
    ("Wangda Zuo (Co-PI)",              "Penn State — AE",  "Energy systems; cyber-physical datasets; industrial AI"),
    ("Kathleen (Kathy) Hill (Co-PI)",   "Penn State CSATS", "K-12 teacher PD; rural school partnerships; implementation"),
    ("Amber Cesare",                    "CSATS",            "Curriculum integration; complexity dial tuning; field liaison"),
    ("Jeff Remington",                  "CSATS",            "School outreach; district relationships; NEIC facilitation"),
]
col_w   = [4.2, 2.8, 6.1]
col_x   = [0.25, 4.55, 7.45]
row_h   = 0.72
top0    = 1.45

# header row
for ci, (label, cx, cw) in enumerate(zip(["Name / Role", "Institution", "Expertise"], col_x, col_w)):
    rect(s, cx, top0, cw - 0.05, 0.38, fill_rgb=PSU_BLUE)
    txbox(s, label, cx + 0.06, top0 + 0.04, cw - 0.15, 0.32,
          size=11, bold=True, color=PSU_WHITE)

for ri, (name, inst, exp) in enumerate(rows):
    bg = LIGHT_BG if ri % 2 == 0 else PSU_WHITE
    y = top0 + 0.38 + ri * row_h
    for ci, (val, cx, cw) in enumerate(zip([name, inst, exp], col_x, col_w)):
        rect(s, cx, y, cw - 0.05, row_h - 0.04, fill_rgb=bg)
        txbox(s, val, cx + 0.06, y + 0.05, cw - 0.18, row_h - 0.1,
              size=10, color=TEXT_DARK)


# ── Slide 3 — Network Partners ───────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
header_bar(s, "Network Partners")
slide_number(s, 3)

# Left column — schools
txbox(s, "K-12 School Partners — 7 Rural PA Districts", 0.3, 1.45, 7.5, 0.4,
      size=13, bold=True, color=PSU_BLUE)

schools = [
    "Keystone Central: Bucktail MS/HS  (100% econ. disadvantaged)",
    "Keystone Central: Central Mountain HS  (64% econ. disadvantaged)",
    "Juniata Valley HS  (56% econ. disadvantaged)",
    "Juniata HS  (42% econ. disadvantaged)",
    "Hollidaysburg Senior HS  (37% econ. disadvantaged)",
    "Lewisburg Area HS  (27% econ. disadvantaged)",
    "Philipsburg-Osceola Senior HS  (46% econ. disadvantaged)",
    "Jersey Shore Senior HS  (47% econ. disadvantaged)",
]
for i, sch in enumerate(schools):
    y = 1.9 + i * 0.55
    rect(s, 0.3, y, 7.2, 0.48, fill_rgb=PSU_WHITE)
    txbox(s, "•  " + sch, 0.45, y + 0.06, 7.0, 0.38, size=11, color=TEXT_DARK)

# Right column — data + infra partners
rect(s, 7.9, 1.35, 5.1, 5.8, fill_rgb=PSU_WHITE)
txbox(s, "Data & Infrastructure Partners", 8.0, 1.45, 4.8, 0.4,
      size=13, bold=True, color=PSU_BLUE)

partners = [
    ("2DCC-MIP (Penn State)", "Materials science AFM datasets — source for Composable Courseware industrial context"),
    ("Penn State ICDS", "Compute + storage infrastructure for telemetry backend and dataset archiving"),
]
py = 1.95
for name, desc in partners:
    txbox(s, name, 8.0, py, 4.7, 0.35, size=11, bold=True, color=PSU_BLUE)
    txbox(s, desc, 8.0, py + 0.33, 4.7, 0.55, size=10, color=MID_GRAY)
    py += 1.05


# ── Slide 4 — The K-12 Math Problem ─────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
header_bar(s, "The K-12 Math Problem We're Addressing")
slide_number(s, 4)

# Big quote / hook
rect(s, 0.3, 1.4, 12.7, 1.0, fill_rgb=PSU_BLUE)
txbox(s, "Real-world data is messy.  Textbook data isn't.",
      0.5, 1.5, 12.3, 0.8,
      size=22, bold=True, color=PSU_WHITE, align=PP_ALIGN.CENTER)

body = (
    "Students learn mathematics with sanitized, curated datasets — but data in science, engineering, and "
    "industry is high-dimensional, noisy, and complex.  This gap means students don't develop authentic "
    "mathematical modeling competency: the ability to reason through real-world complexity.\n\n"
    "Our partner schools serve students in rural, high-poverty communities where math proficiency rates "
    "range from 12 – 47%.  These students have the most to gain from instruction that connects math to "
    "the world they will enter — and the fewest resources to build that bridge without support."
)
txbox(s, body, 0.4, 2.55, 8.5, 2.8, size=13, color=TEXT_DARK)

# RQ box
rect(s, 9.1, 1.4, 3.9, 4.3, fill_rgb=PSU_WHITE)
txbox(s, "Research Questions", 9.2, 1.5, 3.7, 0.4, size=12, bold=True, color=PSU_BLUE)
txbox(s, "RQ 1\nHow do students develop mathematical modeling competencies when engaging with complex industrial datasets in Algebra and Statistics?",
      9.2, 1.95, 3.65, 1.6, size=10.5, color=TEXT_DARK)
txbox(s, "RQ 2\nWhat infrastructure — data translation frameworks and annotation protocols — is required to study mathematical modeling at scale?",
      9.2, 3.65, 3.65, 1.6, size=10.5, color=TEXT_DARK)

txbox(s, "Target: ~4,500 students | 7 rural PA districts | Years 1–3",
      0.4, 5.5, 8.5, 0.4, size=12, bold=True, color=MID_GRAY)


# ── Slide 5 — What Data We're Building ──────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
header_bar(s, "What Data We're Building")
slide_number(s, 5)

# Two industrial context boxes
for xi, (title, body) in enumerate([
    ("Materials Science Context", "Atomic force microscopy (AFM) data from Penn State's 2D Crystal Consortium (Reinhart)\nStatistical + structural complexity"),
    ("Energy Systems Context", "Campus cooling system and data center time-series (Zuo)\nTemporal + provenance complexity"),
]):
    lx = 0.3 + xi * 6.4
    rect(s, lx, 1.4, 6.1, 1.1, fill_rgb=ACCENT)
    txbox(s, title, lx + 0.1, 1.47, 5.9, 0.35, size=12, bold=True, color=PSU_BLUE)
    txbox(s, body, lx + 0.1, 1.85, 5.9, 0.58, size=10.5, color=TEXT_DARK)

# Five layers
txbox(s, "Five Linked Data Layers per Student Session", 0.3, 2.65, 12.7, 0.38,
      size=14, bold=True, color=PSU_BLUE)

layers = [
    ("1  Process-level telemetry", "Cell executions · code diffs · errors · idle time  (tens of millions of events; JSONL)"),
    ("2  Pre/post assessments",    "Mathematical modeling competency measures before and after each module"),
    ("3  Student work artifacts",  "Completed Jupyter notebooks capturing the full modeling process"),
    ("4  Classroom metadata",      "Teacher fidelity of implementation; session context; site demographics"),
    ("5  Expert annotations",      "Reasoning Pattern Codes · Error Taxonomy · Productive Struggle Indicators  (codebook D6)"),
]
for i, (num, desc) in enumerate(layers):
    y = 3.1 + i * 0.72
    bg = PSU_BLUE if i % 2 == 0 else RGBColor(0x2A, 0x55, 0x9A)
    rect(s, 0.3, y, 2.7, 0.65, fill_rgb=bg)
    txbox(s, num, 0.38, y + 0.08, 2.55, 0.5, size=11, bold=True, color=PSU_WHITE)
    rect(s, 3.05, y, 10.0, 0.65, fill_rgb=PSU_WHITE)
    txbox(s, desc, 3.15, y + 0.08, 9.8, 0.5, size=11, color=TEXT_DARK)

txbox(s, "FAIR-compliant · CC-BY · Published on ICPSR + OSF with DOI · Tutorial notebooks included",
      0.3, 6.75, 12.5, 0.38, size=11, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 6 — Current State ──────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
header_bar(s, "Current State of the Data",
           subtitle="Award received June 22, 2026 — we are in the pre-work setup period (July–December 2026)")
slide_number(s, 6)

statuses = [
    ("Composable Courseware L01–L02",         "Built",                            PSU_BLUE,  "From prior NSF IIS-2123343 — tested in college classrooms, now adapting for K-12"),
    ("School district partnerships",           "Outreach underway",                RGBColor(0xE2,0x97,0x00), "7 districts identified; Kathy Hill's team has existing relationships"),
    ("Python telemetry framework",             "Designed — impl. Jan 2027",        MID_GRAY,  "IPython kernel hooks + FastAPI backend + teacher dashboard"),
    ("Annotation codebook",                    "Drafted — pilot validation needed", MID_GRAY,  "Reasoning Pattern Codes, Error Taxonomy, Productive Struggle Indicators"),
    ("Pilot dataset (~100 students, 2 sites)", "Target: Fall 2027",                MID_GRAY,  "1 Title I + 1 rural school · inter-rater reliability κ ≥ 0.70"),
    ("Full dataset (~4,500 students)",         "Target: End of 2029",              MID_GRAY,  "5 layers · tens of millions events · published ICPSR + OSF"),
]
col_labels = ["Component", "Status", "Notes"]
col_x2 = [0.25, 4.0, 6.3]
col_w2 = [3.65, 2.2, 6.7]
top2 = 1.45

for ci, (label, cx, cw) in enumerate(zip(col_labels, col_x2, col_w2)):
    rect(s, cx, top2, cw - 0.05, 0.38, fill_rgb=PSU_BLUE)
    txbox(s, label, cx + 0.06, top2 + 0.04, cw - 0.15, 0.32,
          size=11, bold=True, color=PSU_WHITE)

row_h2 = 0.75
for ri, (comp, stat, stat_color, note) in enumerate(statuses):
    bg2 = LIGHT_BG if ri % 2 == 0 else PSU_WHITE
    y = top2 + 0.38 + ri * row_h2
    rect(s, col_x2[0], y, col_w2[0] - 0.05, row_h2 - 0.04, fill_rgb=bg2)
    txbox(s, comp, col_x2[0] + 0.06, y + 0.08, col_w2[0] - 0.18, row_h2 - 0.15, size=10, color=TEXT_DARK)
    rect(s, col_x2[1], y, col_w2[1] - 0.05, row_h2 - 0.04, fill_rgb=bg2)
    txbox(s, stat, col_x2[1] + 0.06, y + 0.08, col_w2[1] - 0.18, row_h2 - 0.15, size=10, bold=True, color=stat_color)
    rect(s, col_x2[2], y, col_w2[2] - 0.05, row_h2 - 0.04, fill_rgb=bg2)
    txbox(s, note, col_x2[2] + 0.06, y + 0.08, col_w2[2] - 0.18, row_h2 - 0.15, size=10, color=MID_GRAY)


# ── Slide 7 — How We Draw on Expertise ──────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
header_bar(s, "How We Draw on Researchers, Practitioners, and Data Scientists")
slide_number(s, 7)

cols_exp = [
    ("Researchers", PSU_BLUE, [
        "Yao — cognitive science, science of learning",
        "→ annotation protocols, RQ1 empirical analysis",
        "",
        "Reinhart — materials science, FAIR data",
        "→ dataset originator + Composable Courseware architecture",
        "",
        "Zuo — energy systems engineering",
        "→ second industrial dataset context",
    ]),
    ("Practitioners", RGBColor(0x2A, 0x55, 0x9A), [
        "Hill + Cesare + Remington (CSATS)",
        "→ teacher PD design, school partnerships,",
        "   fidelity monitoring",
        "",
        "4–12 math teachers / year as co-researchers",
        "→ annotating, documenting, shaping tools",
        "",
        "NEIC — monthly virtual sessions, PDSA cycles",
        "→ practitioner feedback loop closes every 30 days",
    ]),
    ("Data Scientists", RGBColor(0x0D, 0x2B, 0x57), [
        "Napolitano — data translation framework",
        "→ complexity dials, scaffolding templates",
        "",
        "2 PhD students (to be hired)",
        "→ telemetry system + ML classifier development",
        "",
        "D8 classifier: proof-of-concept AI on our dataset",
        "→ predicts productive vs. unproductive struggle",
    ]),
]

for ci, (title, color, bullets) in enumerate(cols_exp):
    lx = 0.25 + ci * 4.35
    rect(s, lx, 1.4, 4.1, 0.55, fill_rgb=color)
    txbox(s, title, lx + 0.1, 1.47, 3.9, 0.42, size=15, bold=True, color=PSU_WHITE)
    rect(s, lx, 1.98, 4.1, 5.1, fill_rgb=PSU_WHITE)
    for bi, line in enumerate(bullets):
        txbox(s, line, lx + 0.12, 2.06 + bi * 0.54, 3.85, 0.5,
              size=10.5, color=TEXT_DARK if not line.startswith("→") else MID_GRAY)


# ── Slide 8 — What We Expect to Work Well ───────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
header_bar(s, "What We Expect to Work Well")
slide_number(s, 8)

strengths = [
    ("Not starting from zero",
     "Composable Courseware L01–L02 are built and field-tested from a completed NSF grant (IIS-2123343). Phase I adapts proven architecture for K-12 — not greenfield development."),
    ("Strong school trust already in place",
     "Kathy Hill's CSATS team has existing relationships with all 7 rural districts. Partnership outreach is warm, not cold."),
    ("Interdisciplinary team with shared track record",
     "Napolitano, Reinhart, and Hill collaborated on a completed NSF grant together. We know how to work across AE, MSE, and education."),
    ("Mature telemetry design",
     "IPython kernel hooks for process-level data capture are well-understood from prior work. The engineering design is validated in principle; Year 1 is implementation, not invention."),
]

for i, (title, desc) in enumerate(strengths):
    y = 1.45 + i * 1.42
    rect(s, 0.3, y, 0.55, 1.25, fill_rgb=PSU_BLUE)
    txbox(s, str(i+1), 0.3, y + 0.35, 0.55, 0.55,
          size=24, bold=True, color=PSU_WHITE, align=PP_ALIGN.CENTER)
    rect(s, 0.9, y, 12.1, 1.25, fill_rgb=PSU_WHITE)
    txbox(s, title, 1.05, y + 0.08, 11.8, 0.38, size=13, bold=True, color=PSU_BLUE)
    txbox(s, desc, 1.05, y + 0.5, 11.8, 0.7, size=11, color=TEXT_DARK)


# ── Slide 9 — Biggest Challenges ────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BG)
header_bar(s, "Biggest Challenges — and How a Collaboratory Helps")
slide_number(s, 9)

challenges = [
    ("Data translation is genuinely hard",
     "Making industrial data pedagogically accessible for Algebra/Statistics requires expertise that rarely sits in one person. Our 'complexity dials' are the mechanism — calibrating them across 7 districts, multiple teachers, and diverse students is an open research problem."),
    ("Annotation at scale requires reproducible standards",
     "Labeling student reasoning (What counts as 'productive struggle'?) must be precise enough for non-expert coders to reach κ ≥ 0.70 across thousands of sessions. Getting there — and sustaining it — is a major infrastructure challenge."),
    ("The logistics chain is tight",
     "Every milestone depends on IRB approval → partnership agreements → hiring → training. This chain takes months. Anything that slows the front end ripples through Year 1."),
]

for i, (title, desc) in enumerate(challenges):
    y = 1.42 + i * 1.52
    rect(s, 0.3, y, 0.55, 1.38, fill_rgb=RGBColor(0xC0, 0x39, 0x2B))
    txbox(s, str(i+1), 0.3, y + 0.38, 0.55, 0.6,
          size=24, bold=True, color=PSU_WHITE, align=PP_ALIGN.CENTER)
    rect(s, 0.9, y, 12.1, 1.38, fill_rgb=PSU_WHITE)
    txbox(s, title, 1.05, y + 0.08, 11.8, 0.38, size=13, bold=True, color=RGBColor(0xC0,0x39,0x2B))
    txbox(s, desc, 1.05, y + 0.5, 11.8, 0.82, size=11, color=TEXT_DARK)

# Collaboratory box
rect(s, 0.3, 6.0, 12.7, 1.15, fill_rgb=ACCENT)
txbox(s, "Where a collaboratory helps most:", 0.45, 6.05, 12.3, 0.35, size=11, bold=True, color=PSU_BLUE)
txbox(s,
      "Shared annotation standards → pool coding labor across teams  ·  "
      "Cross-team data translation templates → faster complexity dial development  ·  "
      "Shared IRB/data governance models → reduce every team's setup burden",
      0.45, 6.4, 12.3, 0.65, size=10.5, color=TEXT_DARK)


# ── Slide 10 — Magic Wand ────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill_rgb=PSU_BLUE)
slide_number(s, 10)

txbox(s, "✦  Magic Wand", 0.5, 0.4, 12.3, 0.7,
      size=30, bold=True, color=PSU_WHITE)
txbox(s, "If you could solve your biggest challenge instantly — what would that look like?",
      0.5, 1.15, 12.3, 0.5, size=15, color=ACCENT)

rect(s, 0.5, 1.85, 12.3, 3.5, fill_rgb=RGBColor(0x0D, 0x2B, 0x57))
txbox(s,
      "\"If we could wave a magic wand: every CAMEL team would adopt the same process-level "
      "telemetry schema from day one — a shared event log standard — so that student reasoning "
      "data across materials science, climate, and energy contexts could be directly compared.\n\n"
      "Right now we're each building data infrastructure in parallel. A collaboratory that "
      "converges on a common telemetry and annotation layer in Phase I would let Phase II ask "
      "questions none of us can ask alone.\"",
      0.75, 2.0, 11.8, 3.1, size=15, color=PSU_WHITE)

rect(s, 0.5, 5.55, 12.3, 0.65, fill_rgb=FLAG_GOLD)
txbox(s, "⚑  PLACEHOLDER — Replace with Becca's voice before presenting",
      0.65, 5.6, 12.0, 0.52, size=12, bold=True, color=TEXT_DARK)

txbox(s, "NSF CAMEL Phase I  ·  Award #2621173  ·  Penn State University  ·  nap@psu.edu",
      0.5, 6.9, 12.3, 0.4, size=10, color=ACCENT, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(__file__), "2026-07-28-project-showcase.pptx")
prs.save(out)
print(f"Saved: {out}")
